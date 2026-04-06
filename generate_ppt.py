
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK = RGBColor(0x0D, 0x0D, 0x1A)
BG_CARD = RGBColor(0x16, 0x16, 0x2B)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
ACCENT_ORANGE = RGBColor(0xFF, 0xA5, 0x02)
ACCENT_GREEN = RGBColor(0x2E, 0xD5, 0x73)
ACCENT_RED = RGBColor(0xFF, 0x47, 0x57)
ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)

def add_bg(slide):
    """Add dark background to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_shape_bg(slide, left, top, width, height, color=BG_CARD, alpha=None):
    """Add a rounded rectangle card background"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_CYAN):
    """Add a thin accent line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_table(slide, left, top, width, height, rows_data, col_widths=None):
    """Add a styled table"""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Calibri'
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.alignment = PP_ALIGN.LEFT

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
            elif row_idx % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            else:
                cell_fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x28)

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    return table_shape

# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

# Accent lines
add_accent_line(slide, Inches(1), Inches(2.6), Inches(11.3), ACCENT_CYAN)
add_accent_line(slide, Inches(1), Inches(2.68), Inches(11.3), PURPLE)

# Title
add_text_box(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.5),
             "🏁 BHAUKAL RACING", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(0.6),
             "A 2D Physics-Based Multiplayer Hill Climb Racing Game",
             font_size=22, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.CENTER)

# Details card
add_shape_bg(slide, Inches(3.5), Inches(3.2), Inches(6.3), Inches(2.8), BG_CARD)
details = [
    "👤  Presented by: [Your Name]",
    "🎓  Course: [Your Course / Subject]",
    "🏫  Institution: [Your College Name]",
    "📅  Date: March 2026",
    "💻  Platform: Web Browser (Desktop)"
]
add_bullet_list(slide, Inches(4.0), Inches(3.4), Inches(5.5), Inches(2.5),
                details, font_size=18, color=LIGHT_GRAY, spacing=Pt(10))

# Bottom tagline
add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "React.js  •  Phaser 3  •  Node.js  •  Socket.IO  •  Matter.js",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Speaker notes
slide.notes_slide.notes_text_frame.text = (
    "Good morning everyone. Today I'll be presenting Bhaukal Racing — "
    "a 2D physics-based hill climbing racing game built entirely for the web browser. "
    "It features solo play, AI bot opponents, and real-time online multiplayer."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Problem Statement", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_RED)

# Problem card
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.8), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
             "❌ The Problem", font_size=22, color=ACCENT_RED, bold=True)
problems = [
    "•  Most racing games require heavy downloads (500MB–5GB)",
    "•  No multiplayer hill-climb racing in the browser",
    "•  Existing browser games have poor physics quality",
    "•  Casual gamers need quick, instant-play experiences",
    "•  No competitive AI opponents in web racing games",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(3.5),
                problems, font_size=16, color=LIGHT_GRAY, spacing=Pt(12))

# Solution card
add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.8), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.5),
             "✅ Our Solution", font_size=22, color=ACCENT_GREEN, bold=True)
solutions = [
    "•  Zero-install — runs in any modern browser",
    "•  Real physics engine — realistic vehicle behavior",
    "•  3 game modes — Solo, Bot Race, Multiplayer",
    "•  6 unique maps with distinct environments",
    "•  Premium visuals with AAA-style UI design",
]
add_bullet_list(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(3.5),
                solutions, font_size=16, color=LIGHT_GRAY, spacing=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "The core problem is that racing games require large downloads, and browser-based "
    "alternatives lack physics depth. There's no web game combining hill-climb physics "
    "with real-time multiplayer. Bhaukal Racing solves this with zero-install, real physics, "
    "and competitive gameplay."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: OBJECTIVES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Project Objectives", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_CYAN)

objectives = [
    ("🎯", "Build a physics-driven 2D racing game with realistic vehicle handling"),
    ("🎮", "Implement 3 game modes: Solo, AI Bots, Online Multiplayer"),
    ("🗺️", "Design 6 unique maps with distinct environments and physics"),
    ("🤖", "Create competitive AI bots with human-like racing behavior"),
    ("🌐", "Develop real-time multiplayer system using WebSockets"),
    ("⚡", "Implement Nitro Boost, Coins, Fuel, and Trick reward systems"),
    ("🎨", "Deliver a premium, responsive UI with modern design aesthetics"),
    ("📊", "Ensure 60 FPS performance with optimized rendering"),
]

for i, (icon, text) in enumerate(objectives):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 1.35)
    add_shape_bg(slide, x, y, Inches(5.8), Inches(1.1), BG_CARD)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.7),
                 icon, font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.2), Inches(4.8), Inches(0.7),
                 text, font_size=15, color=LIGHT_GRAY)

slide.notes_slide.notes_text_frame.text = (
    "Our objectives were clear — build a complete racing experience with realistic physics, "
    "multiple game modes, diverse maps, and polished UI. Every objective here has been fully implemented."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: KEY FEATURES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Key Features", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_YELLOW)

features_data = [
    ["Category", "Feature"],
    ["🚗 Vehicles", "Car & Bike — each with unique physics and handling"],
    ["🗺️ Maps", "6 maps: Meadows, Desert, Volcano, Highway, Moon, Rooftop"],
    ["🤖 AI Bots", "3 opponents with Easy / Medium / Hard difficulty tiers"],
    ["🌐 Multiplayer", "Real-time online racing with room codes (up to 4 players)"],
    ["⚡ Nitro Boost", "2-use boost system with 45s cooldown and UI feedback"],
    ["🪙 Coins", "5-tier value system (5 / 10 / 25 / 50 / 100) with visual effects"],
    ["⛽ Fuel", "Drain + refuel mechanic with on-track fuel cans"],
    ["🔄 Tricks", "Backflips, multi-flips, wheelies — earn bonus coins"],
    ["🏆 Leaderboard", "Live position tracking and end-of-race rankings"],
    ["🎵 Audio", "Procedurally generated engine, nitro, and pickup sounds"],
]
add_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5), features_data,
          col_widths=[Inches(2.2), Inches(9.5)])

slide.notes_slide.notes_text_frame.text = (
    "Here are the key features — 2 vehicle types, 6 maps with unique physics, "
    "a nitro boost system, 5-tier coin collection, real-time multiplayer, and skill-based tricks."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: TECHNOLOGY STACK
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Technology Stack", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), PURPLE)

# Frontend card
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(3.2), Inches(0.5),
             "🖥️ Frontend", font_size=22, color=ACCENT_CYAN, bold=True)
fe_items = [
    "•  React.js 19 — UI Framework",
    "•  Phaser 3.90 — Game Engine",
    "•  Matter.js — Physics Engine",
    "•  Vite 8 — Build Tool + HMR",
    "•  JavaScript ES6+ — Language",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(3.2), Inches(3.5),
                fe_items, font_size=15, color=LIGHT_GRAY, spacing=Pt(14))

# Backend card
add_shape_bg(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(5.1), Inches(1.6), Inches(3.2), Inches(0.5),
             "⚙️ Backend", font_size=22, color=ACCENT_GREEN, bold=True)
be_items = [
    "•  Node.js — Server Runtime",
    "•  Express 5 — HTTP Server",
    "•  Socket.IO 4 — WebSockets",
    "•  CORS — Cross-Origin Handling",
    "•  Concurrently — Dev Runner",
]
add_bullet_list(slide, Inches(5.1), Inches(2.2), Inches(3.2), Inches(3.5),
                be_items, font_size=15, color=LIGHT_GRAY, spacing=Pt(14))

# Architecture card
add_shape_bg(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(9.1), Inches(1.6), Inches(3.2), Inches(0.5),
             "🏗️ Architecture", font_size=22, color=ACCENT_ORANGE, bold=True)
arch_items = [
    "•  Client-Server Model",
    "•  Custom EventBus System",
    "•  Component-Based UI",
    "•  Modular Game Entities",
    "•  Real-Time WebSocket Sync",
]
add_bullet_list(slide, Inches(9.1), Inches(2.2), Inches(3.2), Inches(3.5),
                arch_items, font_size=15, color=LIGHT_GRAY, spacing=Pt(14))

slide.notes_slide.notes_text_frame.text = (
    "Frontend uses React for UI and Phaser 3 as the game engine with Matter.js for physics. "
    "Backend runs on Node.js with Express and Socket.IO. The architecture follows a client-server "
    "model with WebSocket-based real-time communication."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "System Architecture", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_CYAN)

# CLIENT box
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(3.0),
             RGBColor(0x0F, 0x1A, 0x2E))
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(3), Inches(0.4),
             "CLIENT (Browser)", font_size=18, color=ACCENT_CYAN, bold=True)

# React UI sub-box
add_shape_bg(slide, Inches(1.0), Inches(2.1), Inches(3.8), Inches(2.1), DARK_GRAY)
add_text_box(slide, Inches(1.2), Inches(2.15), Inches(3.5), Inches(0.3),
             "React UI Layer", font_size=13, color=ACCENT_CYAN, bold=True)
ui_items = ["StartScreen.jsx", "GameScreen.jsx (HUD)", "ResultScreen.jsx", "LobbyScreen.jsx"]
add_bullet_list(slide, Inches(1.2), Inches(2.5), Inches(3.4), Inches(1.5),
                ui_items, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# Phaser sub-box
add_shape_bg(slide, Inches(5.1), Inches(2.1), Inches(3.9), Inches(2.1), DARK_GRAY)
add_text_box(slide, Inches(5.3), Inches(2.15), Inches(3.5), Inches(0.3),
             "Phaser 3 Game Engine", font_size=13, color=ACCENT_GREEN, bold=True)
phaser_items = ["GameScene.js (Main Loop)", "Car.js / Bike.js (Vehicles)", "BotManager.js (AI)", "NitroSystem.js / TerrainGen"]
add_bullet_list(slide, Inches(5.3), Inches(2.5), Inches(3.5), Inches(1.5),
                phaser_items, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# Arrow
add_text_box(slide, Inches(4.0), Inches(4.6), Inches(2), Inches(0.5),
             "⬇  WebSocket  ⬇", font_size=14, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# SERVER box
add_shape_bg(slide, Inches(0.8), Inches(5.1), Inches(8.5), Inches(1.8),
             RGBColor(0x1A, 0x0F, 0x0F))
add_text_box(slide, Inches(1.0), Inches(5.15), Inches(3), Inches(0.4),
             "SERVER (Node.js)", font_size=18, color=ACCENT_ORANGE, bold=True)

server_boxes = [
    ("Express.js\n(HTTP)", Inches(1.0)),
    ("RoomManager\n(Create/Join)", Inches(3.3)),
    ("GameManager\n(Race Sync)", Inches(5.6)),
    ("Socket.IO\n(WebSocket)", Inches(7.5)),
]
for text, x in server_boxes:
    add_shape_bg(slide, x, Inches(5.7), Inches(1.8), Inches(1.0), DARK_GRAY)
    add_text_box(slide, x + Inches(0.1), Inches(5.75), Inches(1.6), Inches(0.9),
                 text, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Game flow on right
add_shape_bg(slide, Inches(9.8), Inches(1.5), Inches(2.8), Inches(5.4), BG_CARD)
add_text_box(slide, Inches(10.0), Inches(1.6), Inches(2.5), Inches(0.4),
             "🔄 Game Flow", font_size=16, color=ACCENT_YELLOW, bold=True)
flow_items = [
    "1. Home Screen",
    "2. Select Mode",
    "3. Select Vehicle",
    "4. Select Map",
    "5. Countdown 3-2-1",
    "6. 🏁 Race!",
    "7. Game End",
    "8. Results Screen",
    "9. Try Again / Exit",
]
add_bullet_list(slide, Inches(10.0), Inches(2.1), Inches(2.5), Inches(4.5),
                flow_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(8))

slide.notes_slide.notes_text_frame.text = (
    "The architecture follows a client-server model. The browser runs React for UI and Phaser for the game engine. "
    "They communicate via a custom EventBus. Socket.IO handles multiplayer sync. "
    "The game flow goes from Home to Mode/Vehicle/Map selection to Racing to Results."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: IMPLEMENTATION DETAILS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Implementation Details", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_GREEN)

impl_sections = [
    ("⚙️ Physics Engine", ACCENT_CYAN, [
        "Chassis + 2 wheels connected by spring constraints",
        "Force-based acceleration with speed governor",
        "Per-map gravity & friction for authentic feel",
    ]),
    ("🗺️ Procedural Terrain", ACCENT_GREEN, [
        "Layered sine-wave terrain generated infinitely",
        "Multi-harmonic: amplitude + frequency + smoothness",
        "Types: hills, craters, buildings, highway ramps",
    ]),
    ("🤖 AI Bot System", ACCENT_ORANGE, [
        "3 difficulty profiles (speed, mistakes, recovery)",
        "Rubber-band catch-up + min speed floor",
        "5 independent nitro triggers per bot",
    ]),
    ("⚡ Nitro System", PURPLE, [
        "State machine: Ready → Active → Gap → Cooldown",
        "2.2x force multiplier + angular stabilization",
        "Same rules for player and bots (fair play)",
    ]),
    ("🪙 Coin System", ACCENT_YELLOW, [
        "Weighted random start (100=5%, 10=35%)",
        "Descending-only with repeat/skip allowed",
        "Animated: bobbing, glow, shine, value labels",
    ]),
]

for i, (title, color, items) in enumerate(impl_sections):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.8)
    add_shape_bg(slide, x, y, Inches(3.8), Inches(2.5), BG_CARD)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_accent_line(slide, x + Inches(0.2), y + Inches(0.55), Inches(2), color)
    bullet_items = ["•  " + item for item in items]
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(1.7),
                    bullet_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(8))

slide.notes_slide.notes_text_frame.text = (
    "Key implementations: Matter.js-based vehicle physics with chassis-wheel constraints. "
    "Procedural terrain using layered sine waves. AI bots with 3 difficulty profiles and rubber-banding. "
    "Nitro is a state machine enforcing identical rules for players and bots. "
    "Coins use weighted spawning in descending-only sequences."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: CHALLENGES FACED
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Challenges Faced & Solutions", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_RED)

challenges_data = [
    ["#", "Challenge", "Solution"],
    ["1", "Vehicle flipping on rough terrain", "Angular velocity limiting + ground stabilization"],
    ["2", "Bots too fast or too slow", "Calibrated speed profiles + min speed floor + rubber banding"],
    ["3", "Coins blending with terrain", "Uniform bright yellow + dark backing + black bold text"],
    ["4", "Nitro causing vehicle instability", "Angular stabilization during boost + smooth ramp"],
    ["5", "Multiplayer sync lag", "Interpolated ghost positions + reduced update frequency"],
    ["6", "HUD elements overlapping", "Flexbox-based flow layout replacing absolute positioning"],
    ["7", "Fuel spawning inside terrain", "Safe-Y validation + slope detection for spawn points"],
    ["8", "Moon map too floaty", "Tuned gravity (0.4x) + increased air control"],
]
add_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5), challenges_data,
          col_widths=[Inches(0.6), Inches(4.5), Inches(6.6)])

slide.notes_slide.notes_text_frame.text = (
    "The biggest challenge was physics stability — vehicles flipping on steep terrain. "
    "Solved with angular velocity capping. Bot balancing required multi-layered speed systems. "
    "UI overlaps were fixed by restructuring the HUD layout."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: RESULTS & OUTPUT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Results & Output", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_GREEN)

# Achievements card
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5.5), Inches(0.4),
             "✅ What We Achieved", font_size=20, color=ACCENT_GREEN, bold=True)
achievements = [
    "✅  Fully playable in any modern web browser",
    "✅  Consistent 60 FPS during gameplay",
    "✅  6 unique maps with distinct physics & visuals",
    "✅  Competitive AI stays within ±10% of player speed",
    "✅  Real-time multiplayer tested with 4 players",
    "✅  Zero-install — no downloads or plugins needed",
    "✅  Premium UI with glassmorphism design",
    "✅  Full audio system (engine, nitro, pickups)",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5.5), Inches(4.0),
                achievements, font_size=15, color=LIGHT_GRAY, spacing=Pt(10))

# Stats card
add_shape_bg(slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(7.5), Inches(1.6), Inches(4.8), Inches(0.4),
             "📊 Project Statistics", font_size=20, color=ACCENT_CYAN, bold=True)

stats_data = [
    ["Metric", "Value"],
    ["Total JS modules", "18 files"],
    ["Lines of code", "8,000+"],
    ["Game engine", "Phaser 3 + Matter.js"],
    ["Maps", "6 (Easy → Hard)"],
    ["Vehicle types", "2 (Car, Bike)"],
    ["AI opponents", "3 bots per race"],
    ["Max multiplayer", "4 players"],
    ["Coin tiers", "5 values (5–100)"],
    ["Frame rate", "60 FPS stable"],
]
add_table(slide, Inches(7.5), Inches(2.2), Inches(4.8), Inches(4.0), stats_data,
          col_widths=[Inches(2.4), Inches(2.4)])

slide.notes_slide.notes_text_frame.text = (
    "The game runs at stable 60 FPS in any browser. All 6 maps are playable with unique physics. "
    "AI bots maintain competitive pacing. Multiplayer tested with 4 concurrent players. "
    "Total codebase is around 8,000 lines across 18 modules."
)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: FUTURE SCOPE & CONCLUSION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Future Scope & Conclusion", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), PURPLE)

# Future scope card
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.0), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
             "🔮 Future Scope", font_size=20, color=PURPLE, bold=True)
future_items = [
    "📱  Mobile touch controls + responsive canvas",
    "🚗  More vehicles — Monster Truck, ATV, Skateboard",
    "🗺️  More maps — Arctic, Underwater, Space Station",
    "🏪  Garage system — vehicle upgrades (engine, tires)",
    "🏆  Global leaderboard with database backend",
    "👤  Player profiles — login, stats, achievements",
    "🎮  Tournament / bracket-based competitive mode",
    "🎨  Vehicle skins — customizable colors & decals",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(3.0),
                future_items, font_size=14, color=LIGHT_GRAY, spacing=Pt(8))

# Conclusion card
add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.0), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.4),
             "📋 Conclusion", font_size=20, color=ACCENT_CYAN, bold=True)
conclusion_items = [
    "Bhaukal Racing delivers a complete, polished",
    "physics-based racing game for the web.",
    "",
    "The project demonstrates skills in:",
    "  ✅  Full-stack development (React + Node.js)",
    "  ✅  Game physics (Matter.js, vehicle dynamics)",
    "  ✅  AI programming (bot behavior + strategy)",
    "  ✅  Real-time networking (Socket.IO)",
    "  ✅  UI/UX design (glassmorphism, responsive)",
]
add_bullet_list(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(3.0),
                conclusion_items, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

# Thank You
add_shape_bg(slide, Inches(3.5), Inches(5.8), Inches(6.3), Inches(1.2),
             RGBColor(0x0A, 0x2A, 0x4A))
add_text_box(slide, Inches(3.5), Inches(5.9), Inches(6.3), Inches(0.6),
             "🏁 Thank You!", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(6.4), Inches(6.3), Inches(0.4),
             "Questions?", font_size=18, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "Future plans include mobile support, more vehicles and maps, a garage upgrade system, "
    "and global leaderboards. In conclusion, Bhaukal Racing demonstrates full-stack development, "
    "game physics, AI programming, real-time networking, and modern UI design. Thank you!"
)

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Bhaukal_Racing_Presentation.pptx")
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
